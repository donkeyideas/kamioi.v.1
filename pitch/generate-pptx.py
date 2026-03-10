"""
Generate Kamioi Pitch Deck as PowerPoint (.pptx)
Run: python generate-pptx.py
Output: kamioi-pitch-deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
BG = RGBColor(0x0A, 0x0F, 0x1C)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
DIM = RGBColor(0x94, 0xA3, 0xB8)
MUTED = RGBColor(0x64, 0x74, 0x8B)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
TEAL = RGBColor(0x06, 0xB6, 0xD4)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
PINK = RGBColor(0xEC, 0x48, 0x99)
CARD_BG = RGBColor(0x12, 0x17, 0x2A)
CARD_BORDER = RGBColor(0x1E, 0x29, 0x3B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1)
    shape.rotation = 0.0
    return shape

def add_gradient_bar(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_multiline(slide, left, top, width, height, lines, default_size=16, default_color=DIM):
    """lines = list of (text, size, color, bold, align)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        text = line[0]
        size = line[1] if len(line) > 1 else default_size
        color = line[2] if len(line) > 2 else default_color
        bold = line[3] if len(line) > 3 else False
        align = line[4] if len(line) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Calibri'
        p.alignment = align
        p.space_after = Pt(4)
    return txBox

# ═══════════════════════════════════════════
# SLIDE 1: Cover
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)

# Pre-Seed pill
pill = add_card(slide, 5.5, 1.2, 2.3, 0.45)
add_text(slide, 5.5, 1.22, 2.3, 0.45, 'Pre-Seed  ·  2026', size=13, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, 1.5, 2.0, 10.3, 1.5, 'Turn Every Purchase\nInto an Investment', size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 2.5, 3.8, 8.3, 1.0,
    'AI-powered micro-investing that automatically matches your\nspending to the stocks of the brands you buy.',
    size=20, color=DIM, align=PP_ALIGN.CENTER)

# Stats row
for i, (val, label) in enumerate([('$1–5', 'Per purchase invested'), ('AI', 'Merchant matching'), ('Live', 'Platform built')]):
    x = 3.0 + i * 2.8
    add_text(slide, x, 5.0, 2.0, 0.6, val, size=36, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, 5.55, 2.0, 0.4, label, size=12, color=MUTED, align=PP_ALIGN.CENTER)

add_text(slide, 2, 6.6, 9.3, 0.4, 'www.kamioi.com  ·  Alain Beltran, Founder  ·  Patent Pending', size=14, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)
add_text(slide, 0.8, 0.5, 3, 0.35, 'THE PROBLEM', size=11, color=MUTED, bold=True)
add_text(slide, 0.8, 0.9, 10, 0.8, 'Most people want to invest but never start', size=36, color=WHITE, bold=True)

# Stat cards
stats = [
    ('56%', 'of Americans can\'t cover a $1,000 emergency expense', RED),
    ('70M+', 'Americans don\'t invest at all — intimidated by complex platforms', AMBER),
    ('$0', 'connection between daily spending and wealth building', PURPLE),
]
for i, (val, desc, color) in enumerate(stats):
    x = 0.8 + i * 4.0
    add_card(slide, x, 2.1, 3.7, 2.4)
    add_text(slide, x + 0.3, 2.3, 3.1, 0.7, val, size=42, color=color, bold=True)
    add_text(slide, x + 0.3, 3.1, 3.1, 1.2, desc, size=14, color=DIM)

# Bottom insight
add_card(slide, 0.8, 4.9, 11.7, 1.2)
add_text(slide, 1.2, 5.0, 11.0, 1.0,
    'The gap: People interact with public companies every day — buying coffee at Starbucks,\nshoes from Nike, groceries from Costco — yet they never participate in the upside.\nInvesting is completely disconnected from daily life.',
    size=15, color=DIM)

# ═══════════════════════════════════════════
# SLIDE 3: The Solution
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)
add_text(slide, 0.8, 0.5, 3, 0.35, 'THE SOLUTION', size=11, color=MUTED, bold=True)
add_text(slide, 0.8, 0.9, 10, 0.8, 'Kamioi: spend, match, invest — automatically', size=36, color=WHITE, bold=True)

bullets = [
    'Set a round-up amount ($1, $2, $3+) per purchase — consistent investments, not spare change',
    'AI matches merchants to stocks. Buy Nike? Kamioi invests in $NKE. Starbucks? You own $SBUX.',
    'Receipt scanning breaks multi-item purchases into weighted investments across every brand.',
    'Portfolio grows with every swipe. No decisions, no research — your spending IS your strategy.',
]
for i, b in enumerate(bullets):
    add_text(slide, 1.2, 2.0 + i * 0.65, 6.0, 0.6, f'→  {b}', size=15, color=DIM)

add_text(slide, 1.2, 4.7, 6.0, 0.4, 'PATENT PENDING — Round-up-to-brand-matched-stock investment process', size=11, color=TEAL, bold=True)

# Example cards
examples = [
    ('Coffee  $4.50', '> $2 invested in $SBUX', PURPLE),
    ('Sneakers  $120', '> $2 invested in $NKE', BLUE),
    ('Costco receipt', '> Weighted: $COST, $PG, $KO...', TEAL),
]
for i, (purchase, result, color) in enumerate(examples):
    y = 2.0 + i * 1.2
    add_card(slide, 7.8, y, 4.7, 1.0)
    add_circle(slide, 8.0, y + 0.15, 0.55, color)
    add_text(slide, 8.7, y + 0.08, 3.5, 0.4, purchase, size=13, color=MUTED)
    add_text(slide, 8.7, y + 0.45, 3.5, 0.4, result, size=15, color=WHITE, bold=True)

# ═══════════════════════════════════════════
# SLIDE 4: How It Works
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)
add_text(slide, 0.8, 0.5, 3, 0.35, 'HOW IT WORKS', size=11, color=MUTED, bold=True)
add_text(slide, 0.8, 0.9, 10, 0.8, 'Three steps to autopilot investing', size=36, color=WHITE, bold=True)

steps = [
    ('STEP 1', 'Create Account', 'Sign up in under 2 minutes.\nNo minimums, no experience needed.', PURPLE),
    ('STEP 2', 'Link Your Bank', 'Securely connect your debit card.\nBank-level 256-bit encryption.', BLUE),
    ('STEP 3', 'Set Round-Up', 'Choose $1–$5 per purchase.\nAI handles everything from here.', TEAL),
]
for i, (step, title, desc, color) in enumerate(steps):
    x = 0.8 + i * 4.1
    add_card(slide, x, 2.2, 3.8, 3.2)
    add_circle(slide, x + 1.5, 2.5, 0.7, color)
    add_text(slide, x + 0.3, 3.4, 3.2, 0.3, step, size=11, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.3, 3.7, 3.2, 0.4, title, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.3, 4.2, 3.2, 0.8, desc, size=14, color=DIM, align=PP_ALIGN.CENTER)

add_card(slide, 0.8, 5.8, 11.7, 0.9)
add_text(slide, 1.2, 5.9, 11.0, 0.7,
    'Every purchase after that? Automatically invested into the stock of the brand you bought from.\nYour portfolio grows with every swipe.',
    size=15, color=DIM, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 5: Product
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)
add_text(slide, 0.8, 0.5, 3, 0.35, 'THE PRODUCT', size=11, color=MUTED, bold=True)
add_text(slide, 0.8, 0.9, 10, 0.8, 'Fully built. Ready to launch.', size=36, color=WHITE, bold=True)

features = [
    ('User Dashboard', 'Portfolio tracking, transaction history,\nperformance charts, AI holdings', PURPLE),
    ('AI Engine', 'Multi-provider AI (DeepSeek, Claude,\nOpenAI) for merchant-stock matching', BLUE),
    ('Mobile App', 'React Native / Expo. Portfolio view,\nreceipt capture, bank sync. iOS+Android', TEAL),
    ('Admin Panel', 'User admin, AI cost tracking, CMS,\nblog editor, promo codes, analytics', GREEN),
    ('Receipt Scanner', '3 AI vision providers extract brands\nand create weighted stock allocations', PINK),
    ('Security', '256-bit AES, TOTP 2FA, row-level\nsecurity, SIPC-insured custody', PURPLE),
    ('Blog & SEO', 'Admin-managed blog, auto SEO\nmetadata, structured data, sharing', BLUE),
    ('Live Demo', 'kamioi.com/login\nCode: KAMIOIDEMO', GREEN),
]
for i, (title, desc, color) in enumerate(features):
    row = i // 4
    col = i % 4
    x = 0.8 + col * 3.15
    y = 2.0 + row * 2.3
    add_card(slide, x, y, 2.9, 2.0)
    add_circle(slide, x + 0.2, y + 0.2, 0.45, color)
    add_text(slide, x + 0.2, y + 0.75, 2.5, 0.35, title, size=14, color=WHITE, bold=True)
    add_text(slide, x + 0.2, y + 1.1, 2.5, 0.8, desc, size=11, color=DIM)

# ═══════════════════════════════════════════
# SLIDE 6: Market Opportunity
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)
add_text(slide, 0.8, 0.5, 3, 0.35, 'MARKET OPPORTUNITY', size=11, color=MUTED, bold=True)
add_text(slide, 0.8, 0.9, 10, 0.8, 'A massive, growing market', size=36, color=WHITE, bold=True)

# TAM SAM SOM
for i, (val, label, color, w) in enumerate([
    ('$3.6T', 'Wealth management (2030)', PURPLE, 3.5),
    ('$50B', 'Micro-investing & robo-advisory', BLUE, 3.5),
    ('25%+', 'Annual growth rate', TEAL, 3.5)
]):
    x = 0.8 + i * 4.0
    add_text(slide, x, 2.1, w, 0.6, val, size=40, color=color, bold=True)
    add_text(slide, x, 2.7, w, 0.3, label, size=12, color=MUTED)

# Why now bullets
add_card(slide, 0.8, 3.4, 7.5, 3.2)
add_text(slide, 1.1, 3.5, 7.0, 0.4, 'Why now?', size=18, color=WHITE, bold=True)
why_now = [
    '→  Fractional shares via API (Alpaca, DriveWealth) — infrastructure didn\'t exist 5 years ago',
    '→  Gen Z & Millennials demand passive, automated investing — 73% prefer app-first',
    '→  AI accuracy for merchant matching crossed the reliability threshold',
    '→  Post-COVID wealth gap — 70M+ Americans seeking accessible investment on-ramps',
]
for i, b in enumerate(why_now):
    add_text(slide, 1.3, 4.0 + i * 0.55, 6.8, 0.5, b, size=13, color=DIM)

# TAM funnel
add_card(slide, 8.8, 3.4, 3.7, 3.2)
add_text(slide, 9.0, 3.5, 3.3, 0.35, 'TAM → SAM → SOM', size=11, color=MUTED, bold=True, align=PP_ALIGN.CENTER)

funnel_data = [('$3.6T', 'Total market', PURPLE, 3.0), ('$50B', 'Micro-investing', BLUE, 2.2), ('$500M', 'Year 3 target', TEAL, 1.4)]
for i, (val, label, color, w) in enumerate(funnel_data):
    cx = 8.8 + (3.7 - w) / 2
    y = 4.1 + i * 0.85
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(y), Inches(w), Inches(0.6))
    shape.fill.solid()
    r, g, b = int(str(color)[0:2], 16), int(str(color)[2:4], 16), int(str(color)[4:6], 16)
    shape.fill.fore_color.rgb = RGBColor(r // 4, g // 4, b // 4)
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    add_text(slide, cx + 0.1, y + 0.05, w - 0.2, 0.5, f'{val}  {label}', size=12, color=WHITE, bold=False, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 7: Business Model
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)
add_text(slide, 0.8, 0.5, 3, 0.35, 'BUSINESS MODEL', size=11, color=MUTED, bold=True)
add_text(slide, 0.8, 0.9, 10, 0.8, 'Subscription revenue with clear unit economics', size=36, color=WHITE, bold=True)

# Pricing cards
plans = [
    ('$9.99', '/month', 'Individual', 'Round-up investing, AI matching,\nportfolio tracking, 1 linked card', PURPLE),
    ('$19.99', '/month', 'Family', 'Up to 5 members, shared goals,\nfamily dashboard, 5 linked cards', BLUE),
    ('$49.99', '/month', 'Business', 'Unlimited employees, admin\ndashboard, API, compliance', TEAL),
]
for i, (price, per, name, desc, color) in enumerate(plans):
    x = 0.8 + i * 2.3
    add_card(slide, x, 2.1, 2.1, 2.8)
    add_text(slide, x, 2.3, 2.1, 0.5, price, size=32, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, 2.75, 2.1, 0.3, per, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, x, 3.1, 2.1, 0.3, name, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, 3.5, 2.1, 0.9, desc, size=11, color=DIM, align=PP_ALIGN.CENTER)

# Unit economics
add_card(slide, 0.8, 5.2, 6.1, 1.5)
add_text(slide, 1.1, 5.3, 5.5, 0.3, 'Unit Economics', size=15, color=WHITE, bold=True)
metrics = [('~$0.02', 'Cost per txn'), ('95%+', 'Gross margin'), ('~500', 'Users to breakeven')]
for i, (val, label) in enumerate(metrics):
    x = 1.1 + i * 1.9
    add_text(slide, x, 5.7, 1.7, 0.4, val, size=22, color=GREEN, bold=True)
    add_text(slide, x, 6.1, 1.7, 0.3, label, size=11, color=MUTED)

# Revenue projection
add_card(slide, 7.5, 2.1, 5.0, 4.6)
add_text(slide, 7.8, 2.2, 4.4, 0.4, 'Revenue Projection', size=15, color=WHITE, bold=True)

# Simple bar chart
bars = [('Y1', '$600K', 0.4), ('Y2', '$3.6M', 1.2), ('Y3', '$12M', 2.2), ('Y4', '$36M', 3.2)]
for i, (yr, val, h) in enumerate(bars):
    x = 8.0 + i * 1.1
    bar_bottom = 6.0
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(bar_bottom - h), Inches(0.7), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE if i < 3 else TEAL
    shape.line.fill.background()
    add_text(slide, x - 0.1, bar_bottom - h - 0.35, 0.9, 0.3, val, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, bar_bottom + 0.05, 0.7, 0.25, yr, size=11, color=MUTED, align=PP_ALIGN.CENTER)

add_text(slide, 7.8, 6.35, 4.5, 0.3, 'Based on 5K > 30K > 100K > 300K users, ~$10 avg.', size=10, color=MUTED)

# ═══════════════════════════════════════════
# SLIDE 8: Competition
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)
add_text(slide, 0.8, 0.5, 3, 0.35, 'COMPETITION', size=11, color=MUTED, bold=True)
add_text(slide, 0.8, 0.9, 10, 0.6, 'Why Kamioi wins', size=36, color=WHITE, bold=True)

# Table header
headers = ['Feature', 'Kamioi', 'Acorns', 'Stash', 'Robinhood']
for i, h in enumerate(headers):
    x = 0.8 + i * 2.5
    add_text(slide, x, 1.7, 2.3, 0.35, h, size=11, color=MUTED, bold=True)

# Table rows
rows = [
    ['Brand-matched investing', '✓ AI auto-match', '✗ ETF only', '✗ Manual pick', '✗ Manual pick'],
    ['Round-up investing', '✓ Fixed $1–$5', 'Spare change', '✗ No', '✗ No'],
    ['Receipt scanning', '✓ Multi-AI vision', '✗ No', '✗ No', '✗ No'],
    ['Individual stocks', '✓ Automatic', '✗ ETFs only', '✓ Manual', '✓ Manual'],
    ['Zero decision-making', '✓ Fully passive', '✓ Passive (ETF)', '✗ Active', '✗ Active'],
    ['Target user', 'Passive consumer', 'Beginner saver', 'Beginner investor', 'Active trader'],
]
for r, row in enumerate(rows):
    y = 2.15 + r * 0.55
    # Alternating row bg
    if r % 2 == 0:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y - 0.05), Inches(12.1), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x0F, 0x14, 0x25)
        shape.line.fill.background()
    for c, cell in enumerate(row):
        x = 0.8 + c * 2.5
        color = GREEN if (c == 1 and '✓' in cell) else (MUTED if '✗' in cell else DIM)
        if c == 0:
            color = WHITE
        add_text(slide, x, y, 2.3, 0.4, cell, size=12, color=color, bold=(c == 0))

# Moat card
add_card(slide, 0.8, 6.0, 11.7, 1.0)
add_text(slide, 1.2, 6.1, 11.0, 0.8,
    'Kamioi\'s moat: The AI merchant-matching engine creates a unique data flywheel. Every transaction trains the model.\nReceipt scanning adds a second data channel no competitor has. Brand-loyalty-as-investing is an unclaimed category.\nPatent pending on the core round-up-to-brand-matched-stock process.',
    size=13, color=DIM)

# ═══════════════════════════════════════════
# SLIDE 9: Traction
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)
add_text(slide, 0.8, 0.5, 3, 0.35, 'TRACTION', size=11, color=MUTED, bold=True)
add_text(slide, 0.8, 0.9, 10, 0.8, 'Built from zero to production-ready', size=36, color=WHITE, bold=True)

milestones = [
    ('Platform Architecture', 'React + TypeScript frontend, Supabase\nbackend, PostgreSQL, Edge Functions.\nFull-stack built by solo founder.', PURPLE),
    ('AI Engine Live', 'Multi-provider AI (DeepSeek, Claude,\nOpenAI) for merchant-to-ticker mapping\nwith admin approval & cost tracking.', BLUE),
    ('Receipt Scanner', 'Smart receipt processing with 3 AI\nvision providers. Brand extraction,\nweighted multi-stock allocations.', TEAL),
    ('Mobile App', 'React Native / Expo with portfolio\ntracking, receipt capture, bank sync.\nCross-platform iOS + Android.', GREEN),
]
for i, (title, desc, color) in enumerate(milestones):
    x = 0.8 + i * 3.15
    add_card(slide, x, 2.1, 2.9, 2.8)
    add_circle(slide, x + 1.1, 2.3, 0.55, color)
    add_text(slide, x + 0.2, 3.0, 2.5, 0.35, title, size=14, color=WHITE, bold=True)
    add_text(slide, x + 0.2, 3.4, 2.5, 1.2, desc, size=11, color=DIM)

# Bottom stats
bottom_stats = [
    ('20+', 'Pages shipped', PURPLE),
    ('15+', 'API endpoints', BLUE),
    ('3', 'AI providers', TEAL),
    ('3', 'Dashboard types', GREEN),
    ('1', 'Solo founder', PINK),
]
for i, (val, label, color) in enumerate(bottom_stats):
    x = 0.8 + i * 2.5
    add_card(slide, x, 5.3, 2.2, 1.4)
    add_text(slide, x, 5.5, 2.2, 0.5, val, size=30, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, 6.0, 2.2, 0.3, label, size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 10: Go-to-Market
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)
add_text(slide, 0.8, 0.5, 3, 0.35, 'GO-TO-MARKET', size=11, color=MUTED, bold=True)
add_text(slide, 0.8, 0.9, 10, 0.6, 'Clear path from product to market', size=36, color=WHITE, bold=True)

phases = [
    ('Phase 1 — Now', 'RIA/BD Partnership', 'Secure license/partnership for\nfractional share trading via Alpaca.\nOnly remaining blocker.', PURPLE),
    ('Phase 2 — Month 1–3', 'Beta Launch', 'Invite-only 500 early adopters.\nValidate conversion, retention,\nand AI matching accuracy.', BLUE),
    ('Phase 3 — Month 3–12', 'Growth', 'Content marketing, social media,\nfinance creators, referrals.\nTarget: 5K–10K users Y1.', TEAL),
    ('Phase 4 — Year 2+', 'Scale', 'Family & Business tiers, B2B\npartnerships, white-label.\nSeries A at $5M+ valuation.', GREEN),
]
for i, (phase, title, desc, color) in enumerate(phases):
    y = 1.8 + i * 1.3
    add_circle(slide, 0.9, y + 0.1, 0.35, color)
    add_text(slide, 1.5, y, 2.5, 0.3, phase, size=12, color=color, bold=True)
    add_text(slide, 1.5, y + 0.3, 2.5, 0.3, title, size=16, color=WHITE, bold=True)
    add_text(slide, 4.2, y + 0.05, 3.2, 0.9, desc, size=12, color=DIM)

# Acquisition channels
add_card(slide, 7.8, 1.8, 4.7, 5.0)
add_text(slide, 8.1, 1.9, 4.1, 0.4, 'Acquisition Channels', size=16, color=WHITE, bold=True)
channels = [
    '→  Organic social — "I invested in Nike by buying shoes" is inherently viral',
    '→  Referral program — bonus investment when a friend signs up',
    '→  Finance creators — TikTok/YouTube personal finance influencers',
    '→  SEO & content — educational blog targeting investing queries',
    '→  Receipt sharing — users share "look what I invested in today"',
]
for i, ch in enumerate(channels):
    add_text(slide, 8.1, 2.5 + i * 0.65, 4.2, 0.6, ch, size=12, color=DIM)

add_card(slide, 7.8, 5.5, 4.7, 1.2, fill_color=RGBColor(0x0F, 0x1F, 0x14))
add_text(slide, 8.1, 5.6, 4.2, 1.0,
    'Key insight: Every Kamioi user generates\norganic social proof with every purchase.\n"I invested in $AAPL by buying an iPhone"\nis a marketing message that writes itself.',
    size=12, color=GREEN)

# ═══════════════════════════════════════════
# SLIDE 11: The Ask
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)
add_text(slide, 0.8, 0.5, 3, 0.35, 'THE ASK', size=11, color=MUTED, bold=True)
add_text(slide, 0.8, 0.9, 10, 0.8, 'Raising $250K Pre-Seed', size=36, color=WHITE, bold=True)

# Use of funds
add_card(slide, 0.8, 2.0, 6.5, 4.5)
add_text(slide, 1.2, 2.1, 5.7, 0.4, 'Use of Funds', size=18, color=WHITE, bold=True)

fund_items = [
    ('40%', 'RIA/BD licensing & Alpaca integration', PURPLE, 4.0),
    ('30%', 'Marketing & user acquisition', BLUE, 3.0),
    ('20%', 'Engineering & mobile polish', TEAL, 2.0),
    ('10%', 'Operations & compliance', GREEN, 1.0),
]
for i, (pct, label, color, bar_w) in enumerate(fund_items):
    y = 2.7 + i * 0.9
    add_text(slide, 1.2, y, 0.7, 0.3, pct, size=18, color=color, bold=True)
    add_text(slide, 2.0, y + 0.02, 4.0, 0.3, label, size=14, color=DIM)
    # Progress bar
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(y + 0.35), Inches(bar_w), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# Right side cards
add_card(slide, 7.8, 2.0, 4.7, 1.5)
add_text(slide, 8.1, 2.1, 4.1, 0.35, 'What we need most', size=15, color=WHITE, bold=True)
add_text(slide, 8.1, 2.5, 4.1, 0.9,
    'An RIA or Broker-Dealer partnership to\nprocess fractional share trades through\nAlpaca. Platform is built — this is the final\npiece to go live.',
    size=13, color=DIM)

add_card(slide, 7.8, 3.8, 4.7, 2.0)
add_text(slide, 8.1, 3.9, 4.1, 0.35, 'Milestones this unlocks', size=15, color=WHITE, bold=True)
milestones_text = [
    '→  Secure RIA/BD and launch real trading',
    '→  Onboard 5,000 beta users',
    '→  Reach $5K MRR within 12 months',
    '→  Position for $1M+ Seed round',
]
for i, m in enumerate(milestones_text):
    add_text(slide, 8.1, 4.3 + i * 0.4, 4.1, 0.35, m, size=13, color=DIM)

add_card(slide, 7.8, 6.1, 4.7, 0.9)
add_text(slide, 8.1, 6.15, 4.1, 0.35, 'Terms', size=15, color=WHITE, bold=True)
add_text(slide, 8.1, 6.5, 4.1, 0.35, '$250K SAFE · $2M pre-money cap', size=14, color=TEAL, bold=True)

# ═══════════════════════════════════════════
# SLIDE 12: Founder
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)
add_text(slide, 0.8, 0.5, 3, 0.35, 'THE FOUNDER', size=11, color=MUTED, bold=True)
add_text(slide, 0.8, 0.9, 10, 0.8, 'Built by a builder', size=36, color=WHITE, bold=True)

# Founder card
add_card(slide, 0.8, 2.1, 5.5, 4.5)
add_circle(slide, 1.2, 2.4, 0.9, PURPLE)
add_text(slide, 1.35, 2.55, 0.6, 0.6, 'AB', size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 2.4, 2.4, 3.5, 0.4, 'Alain Beltran', size=22, color=WHITE, bold=True)
add_text(slide, 2.4, 2.8, 3.5, 0.3, 'Founder & CEO', size=13, color=MUTED)

founder_bullets = [
    '→  Full-stack engineer who single-handedly designed, built, and shipped the entire platform',
    '→  Deep expertise: React, TypeScript, Python, Supabase, PostgreSQL, AI/ML, mobile dev',
    '→  Built web app, mobile app, AI engine, receipt scanner, admin dashboard, CMS, and blog',
    '→  Relentless execution — concept to production with 15+ APIs and 20+ pages',
]
for i, b in enumerate(founder_bullets):
    add_text(slide, 1.2, 3.5 + i * 0.6, 4.8, 0.5, b, size=12, color=DIM)

# Why this matters
add_card(slide, 6.8, 2.1, 5.7, 4.5, fill_color=RGBColor(0x10, 0x13, 0x28))
add_text(slide, 7.2, 2.3, 4.9, 0.4, 'Why this matters', size=18, color=WHITE, bold=True)
add_text(slide, 7.2, 2.9, 4.9, 3.0,
    'Most pre-seed startups have a pitch deck\nand a prototype.\n\n'
    'Kamioi has a fully built, production-grade\nplatform — designed, engineered, and\nshipped by one person.\n\n'
    'This demonstrates exceptional technical\nability, product vision, and the drive to\nexecute without external resources.\n\n'
    'Imagine what\'s possible WITH capital,\na team, and the right partnerships.',
    size=14, color=DIM)

# ═══════════════════════════════════════════
# SLIDE 13: Contact / Thank You
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.06)

add_text(slide, 1.5, 1.5, 10.3, 1.2, 'Let\'s build wealth\nfor everyone.', size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 2.5, 3.2, 8.3, 0.6, 'The platform is built. The market is ready.\nWe just need the right partners to launch.', size=20, color=DIM, align=PP_ALIGN.CENTER)

# Contact card
add_card(slide, 4.2, 4.2, 5.0, 2.5)

contacts = [
    ('Website', 'www.kamioi.com', PURPLE),
    ('Live Demo', 'kamioi.com/login  |  code: KAMIOIDEMO', BLUE),
    ('Founder', 'Alain Beltran', TEAL),
]
for i, (label, value, color) in enumerate(contacts):
    y = 4.5 + i * 0.7
    add_text(slide, 4.6, y, 2.0, 0.3, label, size=12, color=MUTED)
    add_text(slide, 6.5, y, 2.5, 0.3, value, size=14, color=WHITE, bold=True)

add_text(slide, 2, 7.0, 9.3, 0.3, 'Kamioi © 2026  ·  Pre-Seed  ·  Raising $250K', size=13, color=MUTED, align=PP_ALIGN.CENTER)

# ── Save ──
out_path = os.path.join(os.path.dirname(__file__), 'kamioi-pitch-deck.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
